import os
from typing import Dict, List, Optional

import fitz  # PyMuPDF
import streamlit as st
from docx import Document
from pptx import Presentation

from src.embedding.pipeline import ingest
from src.integration import build_mindmap_from_entity
from src.integration.dat_adapter import dat_result_to_mindmap_data
from src.interface.layout import (
    close_card,
    open_card,
    render_controls_section,
    render_empty_state,
    render_files_uploaded_state,
    render_graph_empty_state,
    render_header,
    render_processing_state,
    render_render_failed_state,
    render_sidebar_settings,
    render_status_box,
    render_upload_section,
    render_uploaded_files_summary,
    render_results_intro,
)
from src.interface.state import (
    APP_STATE_FILES_UPLOADED,
    APP_STATE_GRAPH_EMPTY,
    APP_STATE_GRAPH_READY,
    APP_STATE_IDLE,
    APP_STATE_PROCESSING,
    APP_STATE_RENDER_FAILED,
    create_default_state,
    reset_processing_outputs,
    set_status,
)
from src.interface.styles import inject_app_css
from src.interface.visualize import render_mindmap_ui
from src.mindmap.structurer import structure_as_study_mindmap
from src.mindmap.study_outputs import generate_study_outputs

st.set_page_config(page_title="MindMap AI", layout="wide")

DATA_DIR = "data"
os.makedirs(DATA_DIR, exist_ok=True)


# =========================
# Session state helpers
# =========================
def init_session_state() -> None:
    if "mindmap_app_state" not in st.session_state:
        st.session_state.mindmap_app_state = create_default_state()


def get_app_state():
    return st.session_state.mindmap_app_state


def save_app_state(state) -> None:
    st.session_state.mindmap_app_state = state


# =========================
# File text extraction
# =========================
def extract_text_from_pdf(uploaded_file) -> str:
    uploaded_file.seek(0)
    pdf_bytes = uploaded_file.read()
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    text = "\n".join(page.get_text() for page in doc)
    doc.close()
    return text


def extract_text_from_pptx(uploaded_file) -> str:
    uploaded_file.seek(0)
    presentation = Presentation(uploaded_file)
    text_runs: List[str] = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_runs.append(shape.text)

    return "\n".join(text_runs)


def extract_text_from_docx(uploaded_file) -> str:
    uploaded_file.seek(0)
    document = Document(uploaded_file)
    return "\n".join(
        [para.text for para in document.paragraphs if para.text.strip()]
    )


def extract_text_from_txt(uploaded_file) -> str:
    uploaded_file.seek(0)
    return uploaded_file.read().decode("utf-8", errors="ignore")


def extract_text_from_file(uploaded_file) -> str:
    filename = uploaded_file.name.lower()

    if filename.endswith(".pdf"):
        return extract_text_from_pdf(uploaded_file)
    if filename.endswith(".pptx"):
        return extract_text_from_pptx(uploaded_file)
    if filename.endswith(".docx"):
        return extract_text_from_docx(uploaded_file)
    if filename.endswith(".txt"):
        return extract_text_from_txt(uploaded_file)

    raise ValueError(f"Unsupported file type: {uploaded_file.name}")


def extract_documents_from_uploaded_files(uploaded_files) -> List[Dict[str, str]]:
    extracted_docs: List[Dict[str, str]] = []

    for uploaded_file in uploaded_files:
        try:
            text = extract_text_from_file(uploaded_file)

            if text.strip():
                extracted_docs.append(
                    {
                        "name": uploaded_file.name,
                        "text": text,
                    }
                )

                base_name = os.path.splitext(uploaded_file.name)[0]
                file_path = os.path.join(DATA_DIR, base_name + ".txt")
                with open(file_path, "w", encoding="utf-8") as f:
                    f.write(text)

        except Exception as e:
            st.error(f"Failed to read {uploaded_file.name}: {e}")

    return extracted_docs


# =========================
# Graph / mindmap pipeline
# =========================
def build_raw_mindmap_result(entity_name: str, graph_store):
    """
    Uses your actual integration function first.
    This may already return adapted mindmap data.
    """
    return build_mindmap_from_entity(
        entity=entity_name.strip(),
        graph_store=graph_store,
    )


def normalize_to_structurable_data(
    raw_result: dict,
    center_topic: str,
    max_neighbors: int,
    prefer_outgoing: bool,
) -> dict:
    """
    Handles both possibilities:
    1. build_mindmap_from_entity returns raw relationship-style data
    2. build_mindmap_from_entity already returns adapted {nodes, edges, center}
    """
    if not raw_result:
        return {"nodes": [], "edges": [], "center": center_topic}

    has_standard_shape = (
        isinstance(raw_result, dict)
        and "nodes" in raw_result
        and "edges" in raw_result
    )

    if has_standard_shape:
        normalized = {
            "nodes": raw_result.get("nodes", []),
            "edges": raw_result.get("edges", []),
            "center": raw_result.get("center") or center_topic,
        }
        return normalized

    return dat_result_to_mindmap_data(
        query_result=raw_result,
        center_topic=center_topic,
        max_neighbors=max_neighbors,
        prefer_outgoing=prefer_outgoing,
    )


# =========================
# UI render helpers
# =========================
def render_uploaded_files_preview(extracted_docs: List[Dict[str, str]]) -> None:
    open_card("Uploaded Files")
    if not extracted_docs:
        st.write("No files uploaded yet.")
    else:
        for doc in extracted_docs:
            st.markdown(f"- {doc['name']}")
    close_card()


def render_source_preview(extracted_docs: List[Dict[str, str]]) -> None:
    open_card("Source Preview")

    if not extracted_docs:
        st.write("No extracted text available.")
        close_card()
        return

    for idx, doc in enumerate(extracted_docs, start=1):
        with st.expander(f"Document {idx}: {doc['name']}"):
            preview_text = doc["text"][:4000]
            st.text_area(
                "Extracted text preview",
                value=preview_text,
                height=240,
                disabled=True,
                key=f"preview_{idx}_{doc['name']}",
            )

    close_card()


def render_study_outputs(study_outputs: Optional[dict]) -> None:
    open_card("Study Outputs")

    if not study_outputs:
        st.write("No study outputs available.")
        close_card()
        return

    tab1, tab2, tab3, tab4, tab5 = st.tabs(
        ["Summary", "Key Topics", "Important Concepts", "Revision Bullets", "Quick Revision"]
    )

    with tab1:
        st.write(study_outputs.get("summary", "No summary available."))

        definitions = study_outputs.get("definitions", [])
        examples = study_outputs.get("examples", [])
        dependencies = study_outputs.get("dependencies", [])
        formulas = study_outputs.get("formulas", [])

        if definitions:
            st.markdown("**Definitions**")
            for item in definitions:
                st.markdown(f"- {item}")

        if examples:
            st.markdown("**Examples**")
            for item in examples:
                st.markdown(f"- {item}")

        if dependencies:
            st.markdown("**Dependencies**")
            for item in dependencies:
                st.markdown(f"- {item}")

        if formulas:
            st.markdown("**Formulas**")
            for item in formulas:
                st.markdown(f"- {item}")

    with tab2:
        for topic in study_outputs.get("key_topics", []):
            st.markdown(f"- {topic}")

    with tab3:
        for concept in study_outputs.get("important_concepts", []):
            st.markdown(f"- {concept}")

    with tab4:
        for bullet in study_outputs.get("revision_bullets", []):
            st.markdown(f"- {bullet}")

    with tab5:
        for line in study_outputs.get("quick_revision", []):
            st.markdown(f"- {line}")

    close_card()


def render_debug_sections(state) -> None:
    with st.expander("Debug: raw graph result"):
        st.json(state.raw_result if state.raw_result else {})

    with st.expander("Debug: adapted graph"):
        st.json(state.adapted_data if state.adapted_data else {})

    with st.expander("Debug: structured mind map"):
        st.json(state.structured_data if state.structured_data else {})


# =========================
# Main app
# =========================
def main() -> None:
    inject_app_css()
    init_session_state()
    state = get_app_state()

    settings = render_sidebar_settings()
    render_header()

    uploaded_files = render_upload_section()

    extracted_docs: List[Dict[str, str]] = []

    if uploaded_files:
        extracted_docs = extract_documents_from_uploaded_files(uploaded_files)
        if extracted_docs:
            state.uploaded_file_name = (
                uploaded_files[0].name if len(uploaded_files) == 1 else f"{len(uploaded_files)} files"
            )
            state.metadata["uploaded_file_names"] = [doc["name"] for doc in extracted_docs]
            state.metadata["documents"] = extracted_docs

            if state.status == APP_STATE_IDLE:
                set_status(state, APP_STATE_FILES_UPLOADED)
                save_app_state(state)

    else:
        if state.status != APP_STATE_IDLE:
            state = create_default_state()
            save_app_state(state)

    documents_from_state = state.metadata.get("documents", [])

    default_topic = ""
    if uploaded_files and len(uploaded_files) == 1:
        default_topic = os.path.splitext(uploaded_files[0].name)[0]
    elif state.uploaded_file_name and state.uploaded_file_name != "files":
        default_topic = ""

    central_topic, generate_clicked = render_controls_section(default_topic=default_topic)

    open_card("Status")

    if state.status == APP_STATE_IDLE:
        render_empty_state()
    elif state.status == APP_STATE_FILES_UPLOADED:
        render_files_uploaded_state(state.uploaded_file_name or "Uploaded file(s)")
    elif state.status == APP_STATE_PROCESSING:
        render_processing_state()
    elif state.status == APP_STATE_GRAPH_EMPTY:
        render_graph_empty_state()
    elif state.status == APP_STATE_RENDER_FAILED:
        render_render_failed_state(state.error_message)
    elif state.status == APP_STATE_GRAPH_READY:
        render_status_box("Mind map generated successfully.")

    close_card()

    if extracted_docs:
        render_uploaded_files_summary([doc["name"] for doc in extracted_docs])

    if generate_clicked:
        if not uploaded_files:
            state = reset_processing_outputs(state)
            set_status(state, APP_STATE_IDLE)
            save_app_state(state)
            st.warning("Please upload at least one file first.")
            return

        if not central_topic.strip():
            st.warning("Please enter a central topic.")
            return

        state = reset_processing_outputs(state)
        set_status(state, APP_STATE_PROCESSING)
        save_app_state(state)

        with st.spinner("Running ingestion and building the study mind map..."):
            try:
                documents = [doc["text"] for doc in extracted_docs]

                graph_store = ingest(documents)
                st.session_state["graph_store"] = graph_store

                raw_result = build_raw_mindmap_result(
                    entity_name=central_topic.strip(),
                    graph_store=graph_store,
                )
                state.raw_result = raw_result

                adapted_data = normalize_to_structurable_data(
                    raw_result=raw_result,
                    center_topic=central_topic.strip(),
                    max_neighbors=settings["max_neighbors"],
                    prefer_outgoing=settings["prefer_outgoing"],
                )
                state.adapted_data = adapted_data

                structured_data = structure_as_study_mindmap(adapted_data)
                state.structured_data = structured_data

                study_outputs = generate_study_outputs(structured_data)
                state.study_outputs = study_outputs

                state.metadata["documents"] = extracted_docs
                state.extracted_text = "\n\n".join(documents)

                if not structured_data.get("nodes") or not structured_data.get("edges"):
                    set_status(state, APP_STATE_GRAPH_EMPTY)
                else:
                    set_status(state, APP_STATE_GRAPH_READY)

            except Exception as e:
                state.error_message = str(e)
                set_status(state, APP_STATE_RENDER_FAILED)

            save_app_state(state)

    state = get_app_state()
    documents_from_state = state.metadata.get("documents", [])

    if state.status == APP_STATE_GRAPH_READY and state.structured_data:
        render_results_intro()
        tab1, tab2, tab3 = st.tabs(["Mind Map", "Study Outputs", "Source Preview"])

        with tab1:
            open_card("Mind Map")
            try:
                render_mindmap_ui(
                    data=state.structured_data,
                    show_debug=settings["show_debug"],
                    prefer_pyvis=settings["prefer_pyvis"],
                )
            except Exception as e:
                state.error_message = str(e)
                set_status(state, APP_STATE_RENDER_FAILED)
                save_app_state(state)
                st.error(f"Render failed: {e}")
            close_card()

        with tab2:
            render_study_outputs(state.study_outputs)

        with tab3:
            if settings["show_source_preview"]:
                render_source_preview(documents_from_state)
            else:
                st.info("Enable source preview from the sidebar to see extracted text.")

    elif state.status == APP_STATE_GRAPH_EMPTY:
        open_card("Results")
        render_graph_empty_state()
        close_card()

        if settings["show_source_preview"]:
            render_source_preview(documents_from_state)

    elif state.status == APP_STATE_RENDER_FAILED:
        open_card("Results")
        render_render_failed_state(state.error_message)
        close_card()

        if settings["show_source_preview"] and documents_from_state:
            render_source_preview(documents_from_state)

    if settings["show_debug"]:
        render_debug_sections(state)


if __name__ == "__main__":
    main()